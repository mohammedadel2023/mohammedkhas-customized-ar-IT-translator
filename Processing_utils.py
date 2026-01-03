import json
import os
import re
import time
from pptx import Presentation
import pdfkit

PATH_WKHTMLTOPDF = r'C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe'

#################################################
#clean part
################################################
def clean_text(text):

    text = text.replace('\x0b', '\n')

    text = re.sub(r'[\u0600-\u06FF\u0750-\u077F\u08A0-\u08FF\uFB50-\uFDFF\uFE70-\uFEFF]', '', text)
    text = re.sub(r'\n?\s*\d+$', '', text)

    lines = text.split('\n')
    cleaned_lines = []
    for line in lines:
        clean_line = ' '.join(line.split())
        if clean_line:
            cleaned_lines.append(clean_line)

    return "\n".join(cleaned_lines)
######################################################
#process_all_pptx
#####################################################

def process_all_pptx(folder_path, output_json_path):
    all_samples = []

    files = [f for f in os.listdir(folder_path)]

    for filename in files:
        file_path = os.path.join(folder_path, filename)
        try:
            prs = Presentation(file_path)
        except Exception as e:
            print(f"Error reading {filename}: {e}")
            continue

        for slide_index, slide in enumerate(prs.slides):
            text_shapes = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text_shapes.append(shape)
            text_shapes.sort(key=lambda x: (x.top, x.left))
            slide_text_content = []

            for shape in text_shapes:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text_content.append(text)

            full_text = "\n".join(slide_text_content)


            cleaned_text = clean_text(full_text)


            if len(cleaned_text.split()) >= 10:
                entry = {
                    "en": cleaned_text,
                    "target": ""
                }
                all_samples.append(entry)
    with open(output_json_path, 'w', encoding='utf-8') as f:
        json.dump(all_samples, f, ensure_ascii=False, indent=4)
#process_all_pptx(r"C:\Users\User\OneDrive\DA350P\files3", r"C:\Users\User\OneDrive\DA350P\training_dataset3.json")

####################################################
#validate_input_file
###################################################

def validate_input_file(file_path, slide_number):
    #Check if file exists
    if not os.path.exists(file_path):
        return 0
    
    
    #Check PowerPoint file extension
    if not file_path.lower().endswith(('.pptx')):
        return 0

    #Check slide number validity
    if not isinstance(slide_number, int) or slide_number <= 0:
        return 0
    #Check slide number bounds
    prs=Presentation(file_path)
    if slide_number > len(prs.slides):
            return 0
    
    
    return 1
################################################
#extract text from specific slides
################################################

def process_single_slide(file_path, slide_number):
    # Convert User Input (1) to Python Index (0)
    slide_index = slide_number - 1
    try:
        prs = Presentation(file_path)
        if slide_index < 0 or slide_index >= len(prs.slides): return 0
        
        slide = prs.slides[slide_index]
        text_shapes = []
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            text_shapes.append(shape)
        text_shapes.sort(key=lambda x: (x.top, x.left))
        
        slide_text_content = []
        for shape in text_shapes:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text: slide_text_content.append(text)
                
        cleaned_text = clean_text("\n".join(slide_text_content))
        
        if len(cleaned_text.split()) < 20: return 0
        
        # Returns a JSON STRING
        return json.dumps({"en": cleaned_text, "target": ""}, ensure_ascii=False)
    except: return 0
########################################################
#HTML_TO_PDF
#######################################################
def render_pdf_from_html_strings(html_text, output_pdf_path):
    if not os.path.exists(PATH_WKHTMLTOPDF):
        print(f"CRITICAL ERROR: wkhtmltopdf not found at {PATH_WKHTMLTOPDF}")
        print("Please install it from https://wkhtmltopdf.org/downloads.html")
        return False

    # Ensure input is a string
    if not isinstance(html_text, str):
        html_text = str(html_text)

    # Wrap in HTML if missing
    if "<html" not in html_text:
        full_html = f"""
        <!DOCTYPE html>
        <html dir="rtl" lang="ar">
        <head>
            <meta charset="UTF-8">
            <style>body {{ font-family: Arial, sans-serif; direction: rtl; text-align: right; }}</style>
        </head>
        <body>{html_text}</body>
        </html>
        """
    else:
        full_html = html_text

    try:
        config = pdfkit.configuration(wkhtmltopdf=PATH_WKHTMLTOPDF)
        
        options = {
            'encoding': "UTF-8",
            'no-outline': None,
            'enable-local-file-access': None
        }
        
        pdfkit.from_string(full_html, output_pdf_path, configuration=config, options=options)
        return True
        
    except Exception as e:
        print(f"PDF Generation Error: {e}")
        return False

######################################################
#unique_output_path
######################################################
def generate_unique_output_path(file_path, slide_number):
    base_name = os.path.splitext(os.path.basename(file_path))[0]
    timestamp = int(time.time())

    os.makedirs("outputs", exist_ok=True)
    filename = f"{base_name}_slide_{slide_number}_{timestamp}.pdf"
    
    return os.path.join("outputs", filename)
